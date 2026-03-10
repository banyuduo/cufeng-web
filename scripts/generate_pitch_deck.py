#!/usr/bin/env python3
"""
簇锋科技 (ToSpike) 融资路演 PPT 生成脚本
文字大纲：完全取自 PPT files/簇锋科技融资路演PPT.pptx
设计风格：完全对齐 PPT files/20260305 公司宣传ppt 2纯净版.svg
运行: python scripts/generate_pitch_deck.py
输出: PPT files/簇锋科技融资路演PPT - 测试.pptx
"""
from ppt_content import FOOTER, SLIDES

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
import os
import tempfile

try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import numpy as np
    HAS_MATPLOTLIB = True
except ImportError:
    HAS_MATPLOTLIB = False

# 配色（完全对齐 20260305 公司宣传ppt 2纯净版.svg）
SVG_BLUE = RGBColor(0x02, 0x3b, 0x99)    # #023B99 主色深蓝
SVG_WHITE = RGBColor(0xff, 0xff, 0xff)    # #FFFFFF 白色
SVG_BG = RGBColor(0xf9, 0xfa, 0xfc)      # #F9FAFC 详情页浅灰背景
SVG_GRAY = RGBColor(0x66, 0x66, 0x66)    # #666666 正文灰 st20
SVG_GRAY_LIGHT = RGBColor(0x80, 0x80, 0x80)  # #808080 次要文字 st6
SVG_ACCENT = RGBColor(0x73, 0xdb, 0xff)  # #73DBFF 强调色 st10
BG_DARK = RGBColor(0x02, 0x3b, 0x99)
BG_DARK_CARD = RGBColor(0x02, 0x3b, 0x99)
BG_LIGHT = RGBColor(0xf9, 0xfa, 0xfc)
TEXT_DARK = RGBColor(0x02, 0x3b, 0x99)
TEXT_BODY = RGBColor(0x66, 0x66, 0x66)
TEXT_MUTED = RGBColor(0x80, 0x80, 0x80)
TEXT_ON_DARK = RGBColor(0xff, 0xff, 0xff)
ACCENT_BLUE = RGBColor(0x73, 0xdb, 0xff)
ACCENT_HIGHLIGHT = RGBColor(0x73, 0xdb, 0xff)
LINE_COLOR = RGBColor(0xe2, 0xe8, 0xf0)
LINE_DARK = RGBColor(0x33, 0x33, 0x33)
WATERMARK = RGBColor(0x02, 0x3b, 0x99)
ACCENT_RED = RGBColor(0xdc, 0x26, 0x26)
FRAME_ACCENT = RGBColor(0x73, 0xdb, 0xff)
# 字体（SVG: FZLTH6K/FZLTH4K 方正兰亭黑, Gotham-Bold/GothamBook, ArialMT）
FONT_TITLE = "FZLanTingHei-R-GBK"         # 方正兰亭黑 Bold 对应 FZLTH6K
FONT_BODY = "FZLanTingHei-R-GBK"         # 方正兰亭黑 Regular 对应 FZLTH4K
FONT_NUM = "Arial"                        # ArialMT
FONT_MONO = "Consolas"
# SVG 资源路径（从 20260305 公司宣传ppt 2纯净版.svg 提取）
PPT_FILES_DIR = os.path.join(os.path.dirname(__file__), "..", "PPT files")
SVG_ASSETS_DIR = os.path.join(PPT_FILES_DIR, "svg_assets")
SVG_LOGO_PATH = os.path.join(SVG_ASSETS_DIR, "logo_banner.jpg")
SVG_PRODUCT_PATH = os.path.join(SVG_ASSETS_DIR, "product_image.jpg")
SVG_TECH_BG_PATH = os.path.join(SVG_ASSETS_DIR, "tech_background.jpg")

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = 100 / 96.0  # 100 像素 ≈ 1.04 英寸
LINE_SPACING = 1.2   # 1.2 倍行高


def apply_title_font(p, size=32):
    """标题：方正兰亭黑 Bold（SVG FZLTH6K st18 约60px）"""
    for fn in (FONT_TITLE, "FZLanTingHeiS-R-GB", "Microsoft YaHei"):
        try:
            p.font.name = fn
            break
        except Exception:
            pass
    p.font.size = Pt(size)
    p.font.bold = True


def apply_body_font(p, size=16):
    """正文：方正兰亭黑 Regular（SVG FZLTH4K st22 约40px）"""
    for fn in (FONT_BODY, "FZLanTingHeiS-R-GB", "Microsoft YaHei"):
        try:
            p.font.name = fn
            break
        except Exception:
            pass
    p.font.size = Pt(size)
    p.font.bold = False
    try:
        p.paragraph_format.line_spacing = LINE_SPACING
    except Exception:
        pass


def add_svg_top_bar(slide, height_inch=0.5):
    """参考 SVG：顶部 #023B99 深蓝条（先添加以置于底层）"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(13.333), Inches(height_inch))
    bar.fill.solid()
    bar.fill.fore_color.rgb = SVG_BLUE
    bar.line.fill.background()


def add_svg_logo(slide, left=0.5, top=0.1, width=2.5, height=None):
    """在幻灯片上添加参考版 LOGO"""
    if os.path.exists(SVG_LOGO_PATH):
        h = height or width * (414 / 736)  # logo 原始比例
        pic = slide.shapes.add_picture(SVG_LOGO_PATH, Inches(left), Inches(top), Inches(width), Inches(h))
        return pic
    return None


def apply_num_font(p):
    """数字使用 Arial"""
    try:
        p.font.name = FONT_NUM
    except Exception:
        pass


def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_line(slide, x1, y1, x2, y2, color=LINE_COLOR, width=0.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)


def add_tech_corner_frame(slide, color=FRAME_ACCENT, corner_len=0.4):
    """Deep-Tech 科技感四角边框：L 形角标"""
    w, h = 13.333, 7.5
    margin = 0.3
    # 左上
    add_line(slide, margin, margin + corner_len, margin, margin, color=color, width=1)
    add_line(slide, margin, margin, margin + corner_len, margin, color=color, width=1)
    # 右上
    add_line(slide, w - margin - corner_len, margin, w - margin, margin, color=color, width=1)
    add_line(slide, w - margin, margin, w - margin, margin + corner_len, color=color, width=1)
    # 左下
    add_line(slide, margin, h - margin - corner_len, margin, h - margin, color=color, width=1)
    add_line(slide, margin, h - margin, margin + corner_len, h - margin, color=color, width=1)
    # 右下
    add_line(slide, w - margin - corner_len, h - margin, w - margin, h - margin, color=color, width=1)
    add_line(slide, w - margin, h - margin - corner_len, w - margin, h - margin, color=color, width=1)


def add_footer(slide, is_dark=False):
    """底部导航（取自原 PPT 文字）"""
    foot = slide.shapes.add_textbox(Inches(MARGIN), Inches(7.15), Inches(12), Inches(0.25))
    tf = foot.text_frame
    p = tf.paragraphs[0]
    p.text = FOOTER
    p.font.size = Pt(8)
    p.font.color.rgb = TEXT_MUTED if not is_dark else RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER


def add_pdf_header_footer(slide, page_label=None, is_dark=False):
    """页眉页脚（兼容旧调用），page_label 用于深色页水印"""
    add_footer(slide, is_dark)
    if page_label and is_dark:
        lbl = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(5), Inches(1.2))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = page_label.replace(" ", "\n")
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WATERMARK


def add_title_box(slide, title, y=1.0, font_size=32, is_dark=False):
    box = slide.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(13.333 - 2 * MARGIN), Inches(0.85))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    apply_title_font(p, font_size)
    p.font.color.rgb = TEXT_ON_DARK if is_dark else TEXT_DARK


def add_content_slide(prs, title, blocks, page_label=None, title_font=32, body_font=16):
    """通用内容页：对齐 SVG 纯净版 — #F9FAFC 背景 + 顶部蓝条 + 方正兰亭黑"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide, page_label)
    add_title_box(slide, title, y=0.55, font_size=title_font)
    add_line(slide, MARGIN, 1.25, 13.333 - MARGIN, 1.25)
    y = 1.9
    for b in blocks:
        if "title" in b:
            hbox = slide.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(13.333 - 2 * MARGIN), Inches(0.6))
            tf = hbox.text_frame
            p = tf.paragraphs[0]
            p.text = b["title"]
            apply_title_font(p, 18)
            p.font.color.rgb = ACCENT_BLUE
            y += 0.7
        for item in b.get("items", []):
            cbox = slide.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(13.333 - 2 * MARGIN), Inches(0.75))
            tf = cbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = ("• " + item) if not item.startswith("•") else item
            apply_body_font(p, body_font)
            p.font.color.rgb = TEXT_BODY
            y += 0.7
        y += 0.2
    return slide


def gen_sp2_sp3_atomic_bonding_img():
    """生成 sp²-sp³ 原子级键合结构图（基于用户 SVG：六边形 sp2 + 圆形 sp3 + 共价键桥接）
    配色符合宣传册深蓝色调 #003366 / #00AEEF"""
    if not HAS_MATPLOTLIB:
        return None
    # 深蓝配色
    STROKE = '#00AEEF'   # 亮蓝描边（与宣传册 #00b4d8 接近）
    FILL_SP3 = '#003366'  # sp3 金刚石深蓝填充
    fig, ax = plt.subplots(figsize=(5, 2.5))
    fig.patch.set_facecolor('#0a1628')
    ax.set_facecolor('#0a1628')
    ax.set_xlim(0, 400)
    ax.set_ylim(0, 200)
    ax.set_aspect('equal')
    ax.axis('off')
    ax.invert_yaxis()
    # sp² 结构（石墨烯六边形）path: M50 100 L70 70 L110 70 L130 100 L110 130 L70 130 Z
    hex_x = [50, 70, 110, 130, 110, 70, 50]
    hex_y = [100, 70, 70, 100, 130, 130, 100]
    ax.plot(hex_x, hex_y, color=STROKE, linewidth=2.5, zorder=2)
    ax.fill(hex_x, hex_y, fill=False)
    # sp³ 结构（金刚石圆形）
    circle = plt.Circle((200, 100), 30, fill=True, facecolor=FILL_SP3, edgecolor=STROKE, linewidth=2)
    ax.add_patch(circle)
    # 化学键合桥梁（虚线）
    ax.plot([130, 170], [100, 100], color=STROKE, linewidth=4, linestyle=(0, (4, 4)), zorder=1)
    # 标签（深蓝背景用浅色文字）
    ax.text(70, 160, 'sp² 导电/导热网', fontsize=11, color='#f5f5f5', ha='center')
    ax.text(200, 160, 'sp³ 刚性骨架', fontsize=11, color='#f5f5f5', ha='center')
    ax.text(150, 85, '共价键桥接', fontsize=10, color=STROKE, ha='center', fontweight='bold')
    plt.tight_layout()
    path = os.path.join(tempfile.gettempdir(), 'tospike_sp2_sp3_atomic.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='#0a1628')
    plt.close()
    return path


def gen_three_stage_rocket_img():
    """三级火箭商业模型图：P0→P1→P2 层级递进，配色体现估值路径"""
    if not HAS_MATPLOTLIB:
        return None
    fig, ax = plt.subplots(figsize=(6, 3.5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 6)
    ax.axis('off')
    # P0 生存 (#003366) -> P1 估值 (#004488) -> P2 护城河 (#0055aa)
    boxes = [
        (2, 4, 2, 1.2, "P0 生存\n线圈/金刚石铜", "#003366"),
        (5, 4, 2, 1.2, "P1 估值\n柔性垫片/不粘锅授权", "#004488"),
        (8, 4, 2, 1.2, "P2 护城河\n固态电池/量子平台", "#0055aa"),
    ]
    for x, y, w, h, text, color in boxes:
        from matplotlib.patches import FancyBboxPatch
        box = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.02", 
                             facecolor=color, edgecolor='#00AEEF', linewidth=2)
        ax.add_patch(box)
        ax.text(x + w/2, y + h/2, text, ha='center', va='center', fontsize=10, 
                color='white', fontweight='bold')
    # 箭头与标签（P0→P1→P2）
    ax.annotate('', xy=(5, 4.6), xytext=(4, 4.6), arrowprops=dict(arrowstyle='->', color='#00AEEF', lw=2))
    ax.text(4.5, 4.9, '产生现金流', fontsize=9, ha='center', color='#1e293b')
    ax.annotate('', xy=(8, 4.6), xytext=(7, 4.6), arrowprops=dict(arrowstyle='->', color='#00AEEF', lw=2))
    ax.text(7.5, 4.9, '技术平移', fontsize=9, ha='center', color='#1e293b')
    plt.tight_layout()
    path = os.path.join(tempfile.gettempdir(), 'tospike_three_stage_rocket.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return path


def gen_blackbox_flowchart_img():
    """黑盒工序流程图：Fabless 模式下 IP 保护，内部核心 vs 外部协作（LR 布局）"""
    if not HAS_MATPLOTLIB:
        return None
    from matplotlib.patches import FancyBboxPatch
    fig, ax = plt.subplots(figsize=(9, 3.2))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 3.5)
    ax.axis('off')
    # 簇锋内部实验室（黑盒区域）
    inner_box = FancyBboxPatch((0.1, 0.8), 3.8, 1.8, boxstyle="round,pad=0.05",
                               facecolor='#fff8f5', edgecolor='#00AEEF', linewidth=2)
    ax.add_patch(inner_box)
    ax.text(2.0, 0.4, '簇锋内部实验室（黑盒区域）', fontsize=10, ha='center', fontweight='bold', color='#003366')
    # 内部流程：原料采购 -> 核心粉末改性 -> 专利包衣处理
    steps_inner = [(0.4, 1.5, "原料\n采购"), (1.4, 1.5, "核心粉末\n改性"), (2.4, 1.5, "专利包衣\n处理")]
    for i, (x, y, t) in enumerate(steps_inner):
        fc = '#ff9966' if i >= 1 else '#e8f4fd'
        ec = '#333333' if i >= 1 else '#94a3b8'
        box = FancyBboxPatch((x, y-0.35), 0.85, 0.7, boxstyle="round,pad=0.02", facecolor=fc, edgecolor=ec, linewidth=1)
        ax.add_patch(box)
        ax.text(x+0.425, y, t, ha='center', va='center', fontsize=8)
        if i < 2:
            ax.annotate('', xy=(x+0.95, y), xytext=(x+0.9, y), arrowprops=dict(arrowstyle='->', color='#333', lw=1))
    # 输出改性母料 ->
    ax.annotate('', xy=(4.1, 1.5), xytext=(3.35, 1.5), arrowprops=dict(arrowstyle='->', color='#00AEEF', lw=2))
    ax.text(3.7, 1.15, '输出改性母料', fontsize=8, ha='center', color='#003366')
    # 外部协作
    outer_box = FancyBboxPatch((4.3, 0.8), 4.0, 1.8, boxstyle="round,pad=0.05",
                               facecolor='#f8fafc', edgecolor='#94a3b8', linewidth=1.5)
    ax.add_patch(outer_box)
    ax.text(6.3, 0.4, '外部协作', fontsize=10, ha='center', color='#64748b')
    steps_outer = [(4.7, 1.5, "外协\n代工厂"), (5.7, 1.5, "真空钎焊\n/成型"), (6.7, 1.5, "精密机\n加工")]
    for i, (x, y, t) in enumerate(steps_outer):
        box = FancyBboxPatch((x, y-0.35), 0.85, 0.7, boxstyle="round,pad=0.02", facecolor='#e2e8f0', edgecolor='#94a3b8', linewidth=1)
        ax.add_patch(box)
        ax.text(x+0.425, y, t, ha='center', va='center', fontsize=8)
        if i < 2:
            ax.annotate('', xy=(x+0.95, y), xytext=(x+0.9, y), arrowprops=dict(arrowstyle='->', color='#333', lw=1))
    # 送回成品 ->
    ax.annotate('', xy=(8.5, 1.5), xytext=(7.6, 1.5), arrowprops=dict(arrowstyle='->', color='#00AEEF', lw=2))
    ax.text(8.05, 1.15, '送回成品', fontsize=8, ha='center', color='#003366')
    # 簇锋内部质检
    box_end = FancyBboxPatch((8.7, 1.15), 1.8, 0.7, boxstyle="round,pad=0.02", facecolor='#e8f4fd', edgecolor='#00AEEF', linewidth=1)
    ax.add_patch(box_end)
    ax.text(9.6, 1.5, '簇锋内部质检', ha='center', va='center', fontsize=9)
    plt.tight_layout()
    path = os.path.join(tempfile.gettempdir(), 'tospike_blackbox_flow.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return path


def gen_phonon_comparison_img():
    """生成声子传输对比图：左灰色混乱散射，右蓝色平滑有序"""
    if not HAS_MATPLOTLIB:
        return None
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(4, 2))
    np.random.seed(42)
    # 左：传统界面，混乱散射
    for _ in range(25):
        x, y = np.random.rand(2) * 0.8 + 0.1
        dx, dy = (np.random.rand(2) - 0.5) * 0.4
        ax1.arrow(x, y, dx, dy, head_width=0.04, head_length=0.03, fc='#94a3b8', ec='#64748b', lw=0.8)
    ax1.set_xlim(0, 1)
    ax1.set_ylim(0, 1)
    ax1.set_aspect('equal')
    ax1.axis('off')
    ax1.set_title('传统界面', fontsize=9, color='#64748b')
    # 右：ToSpike，平滑有序
    for i in range(8):
        y = 0.2 + i * 0.1
        ax2.arrow(0.15, y, 0.7, 0, head_width=0.02, head_length=0.05, fc='#00b4d8', ec='#0e7490', lw=1.2)
    ax2.set_xlim(0, 1)
    ax2.set_ylim(0, 1)
    ax2.set_aspect('equal')
    ax2.axis('off')
    ax2.set_title('ToSpike共价键合', fontsize=9, color='#00b4d8')
    plt.tight_layout()
    path = os.path.join(tempfile.gettempdir(), 'tospike_phonon.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return path


def gen_trend_chart_img():
    """行业趋势/痛点图：AI芯片热流密度 vs 传统材料极限（深色背景 #00AEEF 风格）"""
    if not HAS_MATPLOTLIB:
        return None
    fig, ax = plt.subplots(figsize=(6, 3))
    fig.patch.set_facecolor('#0a1628')
    ax.set_facecolor('#0a1628')
    ax.set_xlim(0, 600)
    ax.set_ylim(0, 300)
    ax.invert_yaxis()
    ax.axis('off')
    # 坐标轴
    ax.plot([50, 550], [250, 250], color='#FFFFFF', linewidth=2)
    ax.plot([50, 50], [250, 50], color='#FFFFFF', linewidth=2)
    # 传统材料天花板 (铜/铝)
    ax.plot([50, 550], [180, 180], color='#FF0000', linewidth=2, linestyle='--')
    ax.text(420, 170, '传统金属散热极限 (300-400W/cm²)', fontsize=10, color='#FF0000')
    # AI芯片热流密度增长曲线 (指数型攀升)
    curve_x = np.linspace(50, 550, 100)
    curve_y = 230 - 190 * ((curve_x - 50) / 500) ** 1.5
    ax.plot(curve_x, curve_y, color='#FFFFFF', linewidth=3)
    ax.text(450, 110, 'AI 芯片功耗需求', fontsize=12, color='#FFFFFF')
    # 簇锋方案覆盖区
    rect = plt.Rectangle((250, 50), 300, 130, facecolor='rgba(0,174,239,0.2)', edgecolor='#00AEEF', linewidth=1)
    ax.add_patch(rect)
    ax.text(400, 70, '簇锋 sp²-sp³ 覆盖区 (>1000W/mK)', fontsize=12, color='#00AEEF', fontweight='bold', ha='center')
    # 时间轴刻度
    ax.text(50, 270, '2020 (A100)', fontsize=10, color='#AAAAAA')
    ax.text(250, 270, '2024 (B200)', fontsize=10, color='#AAAAAA')
    ax.text(450, 270, '2026+ (Next Gen)', fontsize=10, color='#AAAAAA')
    plt.tight_layout()
    path = os.path.join(tempfile.gettempdir(), 'tospike_trend_chart.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='#0a1628')
    plt.close()
    return path


def gen_tech_advantage_radar_img():
    """技术优势对比雷达图：热导率/加工难度/轻量化/成本效率/界面热阻（深色背景 #00AEEF）"""
    if not HAS_MATPLOTLIB:
        return None
    fig, ax = plt.subplots(figsize=(5, 4))
    fig.patch.set_facecolor('#0a1628')
    ax.set_facecolor('#0a1628')
    ax.set_xlim(0, 500)
    ax.set_ylim(0, 400)
    ax.invert_yaxis()
    ax.axis('off')
    # 五边形顶点 (热导率, 加工难度, 轻量化, 成本效率, 界面热阻)
    center = (250, 220)
    vertices = [(250, 50), (430, 180), (360, 360), (140, 360), (70, 180)]
    # 背景网格
    for i, v in enumerate(vertices):
        ax.plot([center[0], v[0]], [center[1], v[1]], color='#444444', linewidth=1)
    inner = [(250, 100), (385, 197), (332, 315), (167, 315), (115, 197)]
    ax.plot([p[0] for p in inner] + [inner[0][0]], [p[1] for p in inner] + [inner[0][1]], color='#333333', linewidth=1)
    ax.plot([p[0] for p in vertices] + [vertices[0][0]], [p[1] for p in vertices] + [vertices[0][1]], color='#444444', linewidth=1)
    # CVD金刚石 (红虚线)
    cvd = [(250, 60), (300, 210), (280, 250), (160, 250), (100, 200), (250, 60)]
    ax.fill([p[0] for p in cvd], [p[1] for p in cvd], facecolor='rgba(255,0,0,0.1)', edgecolor='none')
    ax.plot([p[0] for p in cvd], [p[1] for p in cvd], color='#FF0000', linewidth=2, linestyle='--')
    # 传统金刚石铜 (灰)
    trad = [(250, 150), (340, 190), (310, 290), (190, 290), (160, 200), (250, 150)]
    ax.fill([p[0] for p in trad], [p[1] for p in trad], facecolor='rgba(150,150,150,0.2)', edgecolor='#999999', linewidth=1)
    # 簇锋方案 (亮蓝)
    cufeng = [(250, 70), (410, 185), (350, 350), (150, 350), (85, 185), (250, 70)]
    ax.fill([p[0] for p in cufeng], [p[1] for p in cufeng], facecolor='rgba(0,174,239,0.3)', edgecolor='#00AEEF', linewidth=3)
    # 标签
    ax.text(230, 40, '热导率 (TC)', fontsize=11, color='#FFFFFF')
    ax.text(435, 185, '加工难度\n(近净成形)', fontsize=10, color='#FFFFFF', ha='left')
    ax.text(365, 375, '轻量化', fontsize=11, color='#FFFFFF')
    ax.text(80, 375, '成本效率', fontsize=11, color='#FFFFFF')
    ax.text(10, 185, '界面热阻(低)', fontsize=11, color='#FFFFFF')
    plt.tight_layout()
    path = os.path.join(tempfile.gettempdir(), 'tospike_tech_radar.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='#0a1628')
    plt.close()
    return path


def gen_radar_chart_img():
    """五维雷达图：导热性、减重、可靠性、工艺成本、CTE匹配度"""
    if not HAS_MATPLOTLIB:
        return None
    categories = ['导热性', '减重', '可靠性', '工艺成本', 'CTE匹配度']
    N = len(categories)
    angles = [n / float(N) * 2 * np.pi for n in range(N)]
    angles += angles[:1]
    # ToSpike全碳：五维全满（工艺成本标注中试阶段可控，不夸大，取75）
    tospike = [95, 95, 95, 75, 95]  # 工艺成本75
    tospike += tospike[:1]
    # 金刚石铜：导热性高，其他低
    diamond_cu = [85, 30, 40, 60, 35]
    diamond_cu += diamond_cu[:1]
    # 金刚石碳化硅：各项中等
    diamond_sic = [65, 55, 50, 55, 50]
    diamond_sic += diamond_sic[:1]
    fig, ax = plt.subplots(figsize=(3.5, 3.5), subplot_kw=dict(projection='polar'))
    ax.plot(angles, tospike, 'o-', linewidth=2, color='#00b4d8', label='ToSpike全碳')
    ax.fill(angles, tospike, alpha=0.25, color='#00b4d8')
    ax.plot(angles, diamond_cu, 'o-', linewidth=1.5, color='#64748b', label='金刚石铜')
    ax.fill(angles, diamond_cu, alpha=0.15, color='#94a3b8')
    ax.plot(angles, diamond_sic, 'o-', linewidth=1.5, color='#94a3b8', label='金刚石碳化硅')
    ax.fill(angles, diamond_sic, alpha=0.1, color='#cbd5e1')
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(categories, fontsize=7)
    ax.set_ylim(0, 100)
    ax.legend(loc='upper right', bbox_to_anchor=(1.25, 1.05), fontsize=6)
    plt.tight_layout()
    path = os.path.join(tempfile.gettempdir(), 'tospike_radar.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return path


def set_mono_font(p_or_run, size=10):
    """技术参数使用等宽字体，数字类用 Arial"""
    f = p_or_run.font if hasattr(p_or_run, 'font') else p_or_run
    try:
        f.name = FONT_NUM  # 数字用 Arial
    except Exception:
        try:
            f.name = FONT_MONO
        except Exception:
            pass
    if hasattr(p_or_run, 'font') and hasattr(p_or_run.font, 'size'):
        p_or_run.font.size = Pt(size)


def add_rect_card(slide, left, top, width, height, fill_color=None, no_line=False):
    """添加圆角矩形卡片背景（SVG 纯净版：内容页用白底浅边）"""
    if fill_color is None:
        fill_color = SVG_WHITE  # 内容页默认白底，与 #F9FAFC 背景区分
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if no_line:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = LINE_COLOR
        shape.line.width = Pt(0.25 if fill_color == SVG_WHITE else 0.5)
    shape.adjustments[0] = 0.03
    return shape


def add_cover_slide(prs):
    """Slide 1 封面 — 文字取自原 PPT，设计参考 SVG"""
    c = SLIDES[0]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_DARK)
    add_footer(slide, is_dark=True)
    if os.path.exists(SVG_LOGO_PATH):
        slide.shapes.add_picture(SVG_LOGO_PATH, Inches(1.5), Inches(0.8), Inches(4.5), Inches(2.53))
    title = slide.shapes.add_textbox(Inches(MARGIN), Inches(3.2), Inches(13.333 - 2 * MARGIN), Inches(1.0))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = c["title"]
    apply_title_font(p, 40)
    p.font.color.rgb = TEXT_ON_DARK
    p.alignment = PP_ALIGN.CENTER
    sub = slide.shapes.add_textbox(Inches(MARGIN), Inches(4.2), Inches(13.333 - 2 * MARGIN), Inches(0.5))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = c["subtitle"]
    apply_body_font(p, 16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_HIGHLIGHT
    p.alignment = PP_ALIGN.CENTER
    loc = slide.shapes.add_textbox(Inches(MARGIN), Inches(4.7), Inches(13.333 - 2 * MARGIN), Inches(0.4))
    tf = loc.text_frame
    p = tf.paragraphs[0]
    p.text = c["location"]
    apply_body_font(p, 12)
    p.font.color.rgb = SVG_GRAY_LIGHT
    p.alignment = PP_ALIGN.CENTER
    foot = slide.shapes.add_textbox(Inches(MARGIN), Inches(5.2), Inches(13.333 - 2 * MARGIN), Inches(0.5))
    tf = foot.text_frame
    p = tf.paragraphs[0]
    p.text = c["slogan"]
    apply_body_font(p, 14)
    p.font.color.rgb = SVG_GRAY_LIGHT
    p.alignment = PP_ALIGN.CENTER
    return slide


def add_slide_2_opening(prs):
    """Slide 2 算力革命的隐形瓶颈 — 文字取自原 PPT"""
    c = SLIDES[1]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, c["title"], y=0.55, font_size=26)
    add_line(slide, MARGIN, 1.25, 13.333 - MARGIN, 1.25)
    y = 1.5
    for b in c["blocks"]:
        if "heading" in b:
            hbox = slide.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(12), Inches(0.5))
            tf = hbox.text_frame
            p = tf.paragraphs[0]
            p.text = b["heading"]
            apply_title_font(p, 16)
            p.font.color.rgb = ACCENT_BLUE
            y += 0.55
        for item in b.get("items", []):
            tbox = slide.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(12), Inches(0.6))
            tf = tbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item
            apply_body_font(p, 12)
            p.font.color.rgb = TEXT_BODY
            y += 0.55
        if "footer" in b:
            fbox = slide.shapes.add_textbox(Inches(MARGIN), Inches(y + 0.2), Inches(12), Inches(0.6))
            tf = fbox.text_frame
            p = tf.paragraphs[0]
            p.text = b["footer"]
            apply_body_font(p, 12)
            p.font.bold = True
            p.font.color.rgb = TEXT_DARK
            y += 0.8
        y += 0.15
    if os.path.exists(SVG_PRODUCT_PATH):
        slide.shapes.add_picture(SVG_PRODUCT_PATH, Inches(7.5), Inches(2.0), Inches(5.2), Inches(2.9))
    return slide


def add_slide_3_company(prs):
    """Slide 3 公司概况 — 文字取自原 PPT"""
    c = SLIDES[2]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, c["title"], y=0.55)
    add_line(slide, 0.8, 1.25, 12.5, 1.25)
    y = 1.5
    for b in c["blocks"]:
        hbox = slide.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(12), Inches(0.45))
        tf = hbox.text_frame
        p = tf.paragraphs[0]
        p.text = b["heading"]
        apply_title_font(p, 14)
        p.font.color.rgb = ACCENT_BLUE
        y += 0.5
        for item in b["items"]:
            tbox = slide.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(12), Inches(0.5))
            tf = tbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item
            apply_body_font(p, 10)
            p.font.color.rgb = TEXT_BODY
            y += 0.45
        y += 0.15
    if os.path.exists(SVG_LOGO_PATH):
        slide.shapes.add_picture(SVG_LOGO_PATH, Inches(10.5), Inches(1.6), Inches(2.2), Inches(1.24))
    return slide


def add_slide_4_stages(prs):
    """Slide 4 从工程工具到物理平台 — 文字取自原 PPT"""
    c = SLIDES[3]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, c["title"], y=0.55, font_size=24)
    add_line(slide, 0.8, 1.25, 12.5, 1.25)
    y = 1.6
    stages = c["stages"]
    for i, (stitle, sdesc) in enumerate(stages):
        desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(7.5), Inches(0.9))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sdesc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_BODY
        tag_w, tag_h = 2.8, 0.7
        tag_x = 9.5
        add_rect_card(slide, tag_x, y - 0.05, tag_w, tag_h, fill_color=RGBColor(0xe8, 0xf7, 0xfc), no_line=False)
        tag_box = slide.shapes.add_textbox(Inches(tag_x + 0.1), Inches(y), Inches(tag_w - 0.2), Inches(0.6))
        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        p.text = stitle
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        if i < len(stages) - 1:
            add_line(slide, 0.8, y + 0.55, 12.5, y + 0.55, color=LINE_COLOR)
        y += 1.2
    path_label = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4))
    tf = path_label.text_frame
    p = tf.paragraphs[0]
    p.text = FOOTER
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED
    return slide


def add_slide_5_principle(prs):
    """Slide 5 sp²–sp³原子级共价键合 — 文字取自原 PPT"""
    c = SLIDES[4]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, c["title"], y=0.55, font_size=22)
    add_line(slide, 0.8, 1.25, 12.5, 1.25)
    add_line(slide, 5.5, 1.5, 5.5, 6.2)
    # 左侧：对比图
    add_rect_card(slide, 0.8, 1.5, 4.5, 2.0)
    left = slide.shapes.add_textbox(Inches(0.95), Inches(1.6), Inches(4.1), Inches(0.5))
    tf = left.text_frame
    p = tf.paragraphs[0]
    p.text = c["left"]["title"]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEXT_MUTED
    right = slide.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(4.1), Inches(0.5))
    tf = right.text_frame
    p = tf.paragraphs[0]
    p.text = c["left"]["right_title"]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    sp2_sp3_path = gen_sp2_sp3_atomic_bonding_img()
    if sp2_sp3_path and os.path.exists(sp2_sp3_path):
        slide.shapes.add_picture(sp2_sp3_path, Inches(0.9), Inches(2.5), Inches(4.2), Inches(0.9))
    # 右侧：要点
    hl = slide.shapes.add_textbox(Inches(5.7), Inches(1.5), Inches(6.5), Inches(0.5))
    tf = hl.text_frame
    p = tf.paragraphs[0]
    p.text = c["highlight"]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    y = 2.1
    for pt in c["points"]:
        pbox = slide.shapes.add_textbox(Inches(5.7), Inches(y), Inches(6.5), Inches(0.5))
        tf = pbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = pt
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_BODY
        y += 0.55
    return slide


def add_slide_6_tech_compare(prs):
    """Slide 6 全碳复合材料 vs 现有方案 — 含性能对比表，文字取自原 PPT"""
    c = SLIDES[5]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, c["title"], y=0.55, font_size=18)
    add_line(slide, 0.8, 1.25, 12.5, 1.25)
    sub = slide.shapes.add_textbox(Inches(MARGIN), Inches(1.32), Inches(12), Inches(0.35))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = c["subtitle"]
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED
    # 性能对比表
    tbl = c["table"]
    nrows = len(tbl["rows"]) + 1
    ncols = len(tbl["headers"])
    tshape = slide.shapes.add_table(nrows, ncols, Inches(0.8), Inches(1.7), Inches(11.7), Inches(2.8))
    table = tshape.table
    for col, h in enumerate(tbl["headers"]):
        cell = table.cell(0, col)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
    for row_idx, row in enumerate(tbl["rows"]):
        for col_idx, val in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            p.font.color.rgb = TEXT_BODY
    # 底部说明
    y = 4.6
    for blk in c["blocks"]:
        tbox = slide.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(12), Inches(0.4))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = blk
        apply_body_font(p, 10)
        p.font.color.rgb = TEXT_BODY
        y += 0.45
    fn = slide.shapes.add_textbox(Inches(MARGIN), Inches(5.5), Inches(12), Inches(0.9))
    tf = fn.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = c["footnote"]
    p.font.size = Pt(7)
    p.font.color.rgb = TEXT_MUTED
    return slide


def add_slide_7_diamond_copper(prs):
    """Slide 7 金刚石铜技术纵深 — 文字取自原 PPT"""
    c = SLIDES[6]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, c["title"], y=0.55)
    add_line(slide, 0.8, 1.25, 12.5, 1.25)
    y = 1.6
    for item in c["items"]:
        tbox = slide.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(12), Inches(0.6))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        apply_body_font(p, 12)
        p.font.color.rgb = TEXT_BODY
        y += 0.75
    return slide


def add_slide_8_patent(prs):
    """Slide 8 四层护城河矩阵 — 文字取自原 PPT"""
    c = SLIDES[7]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, c["title"], y=0.55, font_size=20)
    add_line(slide, 0.8, 1.25, 12.5, 1.25)
    cx, cy = 6.5, 4.0
    add_line(slide, cx, 1.5, cx, 6.2)
    add_line(slide, 0.8, cy, 12.5, cy)
    positions = [(0.9, 1.6), (6.7, 1.6), (0.9, 4.1), (6.7, 4.1)]
    for i, (pos, q) in enumerate(zip(positions, c["quadrants"])):
        qbox = slide.shapes.add_textbox(Inches(pos[0]), Inches(pos[1]), Inches(5.4), Inches(2.0))
        tf = qbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = q[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p2 = tf.add_paragraph()
        p2.text = q[1]
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(2)
        p3 = tf.add_paragraph()
        p3.text = q[2]
        p3.font.size = Pt(9)
        p3.font.color.rgb = TEXT_MUTED
        p3.space_before = Pt(2)
    foot = slide.shapes.add_textbox(Inches(MARGIN), Inches(6.4), Inches(12), Inches(0.4))
    tf = foot.text_frame
    p = tf.paragraphs[0]
    p.text = c["footer"]
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    return slide


def add_slide_9_placeholder(prs):
    """Slide 9 空页/图表占位 — 原 PPT 此页无文字"""
    c = SLIDES[8]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    ph = slide.shapes.add_textbox(Inches(4.0), Inches(3.0), Inches(5.3), Inches(1.0))
    tf = ph.text_frame
    p = tf.paragraphs[0]
    p.text = c.get("placeholder", "【插入：图表或图片】")
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER
    return slide


def add_slide_10_rocket(prs):
    """Slide 10 三级火箭模式 — 含主表+金刚石铜产品矩阵，文字取自原 PPT"""
    c = SLIDES[9]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, c["title"], y=0.55, font_size=18)
    add_line(slide, 0.8, 1.25, 12.5, 1.25)
    tbl = c["table"]
    nrows = len(tbl["rows"]) + 1
    ncols = len(tbl["headers"])
    tshape = slide.shapes.add_table(nrows, ncols, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.2))
    table = tshape.table
    for col, h in enumerate(tbl["headers"]):
        cell = table.cell(0, col)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
    for row_idx, row in enumerate(tbl["rows"]):
        for col_idx, val in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(7)
            p.font.color.rgb = TEXT_BODY
    # 金刚石铜产品矩阵子表
    if "subtable" in c:
        st = c["subtable"]
        st_title = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.35))
        tf = st_title.text_frame
        p = tf.paragraphs[0]
        p.text = st.get("title", "金刚石铜产品矩阵")
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        st_nrows = len(st["rows"]) + 1
        st_ncols = len(st["headers"])
        st_shape = slide.shapes.add_table(st_nrows, st_ncols, Inches(0.8), Inches(4.05), Inches(11.7), Inches(1.65))
        st_table = st_shape.table
        for col, h in enumerate(st["headers"]):
            cell = st_table.cell(0, col)
            cell.text = h
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = ACCENT_BLUE
        for row_idx, row in enumerate(st["rows"]):
            for col_idx, val in enumerate(row):
                cell = st_table.cell(row_idx + 1, col_idx)
                cell.text = val
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(7)
                p.font.color.rgb = TEXT_BODY
    foot = slide.shapes.add_textbox(Inches(MARGIN), Inches(5.9), Inches(12), Inches(0.4))
    tf = foot.text_frame
    p = tf.paragraphs[0]
    p.text = c["footer"]
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    return slide


def add_slide_11_team(prs):
    """Slide 11 团队 — 文字取自原 PPT"""
    c = SLIDES[10]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, c["title"], y=0.55, font_size=20)
    add_line(slide, 0.8, 1.25, 12.5, 1.25)
    y = 1.6
    for name, desc in c["members"]:
        nbox = slide.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(12), Inches(0.4))
        tf = nbox.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        y += 0.45
        dbox = slide.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(12), Inches(0.5))
        tf = dbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_BODY
        y += 0.55
    y += 0.2
    gbox = slide.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(12), Inches(0.4))
    tf = gbox.text_frame
    p = tf.paragraphs[0]
    p.text = "治理机制"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    y += 0.45
    for g in c["governance"]:
        gb = slide.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(12), Inches(0.4))
        tf = gb.text_frame
        p = tf.paragraphs[0]
        p.text = g
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_BODY
        y += 0.45
    return slide


def add_slide_12_funding(prs):
    """Slide 12 融资需求 — 文字取自原 PPT"""
    c = SLIDES[11]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, c["title"], y=0.55)
    add_line(slide, 0.8, 1.25, 12.5, 1.25)
    y = 1.7
    cards = [
        ("融资金额", c["amount"]),
        ("出让股份", c["equity"]),
        ("估值逻辑", c["valuation"]),
        ("核心逻辑", c["logic"]),
    ]
    for label, val in cards:
        add_rect_card(slide, 0.8, y, 5.8, 0.9)
        lbl = slide.shapes.add_textbox(Inches(0.95), Inches(y + 0.1), Inches(5.5), Inches(0.35))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        vbox = slide.shapes.add_textbox(Inches(0.95), Inches(y + 0.45), Inches(5.5), Inches(0.4))
        tf = vbox.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_BODY
        y += 1.05
    return slide


def add_slide_13_funding_details(prs):
    """Slide 13 资金用途 — 文字取自原 PPT"""
    c = SLIDES[12]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, c["title"], y=0.55)
    add_line(slide, 0.8, 1.25, 12.5, 1.25)
    y = 1.6
    for cat, items in c["items"]:
        cbox = slide.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(12), Inches(0.4))
        tf = cbox.text_frame
        p = tf.paragraphs[0]
        p.text = cat
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        y += 0.45
        for it in items:
            ibox = slide.shapes.add_textbox(Inches(MARGIN + 0.2), Inches(y), Inches(11.7), Inches(0.4))
            tf = ibox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = it
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_BODY
            y += 0.4
        y += 0.15
    return slide


def add_slide_14_vision(prs):
    """Slide 14 合作愿景 — 文字取自原 PPT"""
    c = SLIDES[13]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, c["title"], y=0.55)
    add_line(slide, 0.8, 1.25, 12.5, 1.25)
    y = 1.6
    for title, items in c["sections"]:
        tbox = slide.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(12), Inches(0.4))
        tf = tbox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        y += 0.5
        for it in items:
            ibox = slide.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(12), Inches(0.5))
            tf = ibox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = it
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_BODY
            y += 0.5
        y += 0.2
    qbox = slide.shapes.add_textbox(Inches(MARGIN), Inches(5.0), Inches(12), Inches(1.2))
    tf = qbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = c["quote"]
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.CENTER
    return slide


def add_slide_5_market_a(prs):
    """AI芯片热管理方案：市场数据 + 客户验证 + 系统级协同示意图（对齐 SVG 纯净版）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, "高性能计算与数据中心：AI 芯片热管理方案", y=0.55, font_size=20)
    add_line(slide, 0.8, 1.25, 12.5, 1.25)
    add_line(slide, 5.5, 1.8, 5.5, 6.2)
    # 左侧：芯片→基板→散热器全链路热管理示意图（标注各环节材料）
    add_rect_card(slide, 0.8, 1.9, 4.5, 2.5)
    # 三个环节框
    boxes = [
        (1.0, 2.1, 1.1, 0.5, "芯片", "Si/封装"),
        (2.3, 2.1, 1.1, 0.5, "基板", "金刚石铜\n全碳复合材料"),
        (3.6, 2.1, 1.1, 0.5, "散热器", "铝/铜+鳍片"),
    ]
    for x, y, w, h, title, mat in boxes:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WATERMARK
        shape.line.color.rgb = ACCENT_BLUE
        shape.line.width = Pt(0.5)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title + "\n" + mat
        p.font.size = Pt(8)
        p.font.color.rgb = TEXT_DARK
        p.alignment = PP_ALIGN.CENTER
    # 箭头（芯片→基板→散热器）
    for (x1, y1, x2, y2) in [(2.1, 2.35, 2.3, 2.35), (3.4, 2.35, 3.6, 2.35)]:
        conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        conn.line.color.rgb = ACCENT_HIGHLIGHT
        conn.line.width = Pt(1)
    # 热流标注
    flow = slide.shapes.add_textbox(Inches(1.5), Inches(2.7), Inches(2.0), Inches(0.4))
    tf = flow.text_frame
    p = tf.paragraphs[0]
    p.text = "热流 → TIM垫片 →"
    p.font.size = Pt(8)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER
    # 右侧：市场数据 + 方案 + 客户验证（压缩高度以留出客户收益表空间）
    right = slide.shapes.add_textbox(Inches(5.7), Inches(1.9), Inches(6.6), Inches(2.4))
    tf = right.text_frame
    tf.word_wrap = True
    blocks = [
        ("市场数据（Yole/IDC）", "先进封装热管理市场快速增长；AI 芯片功率密度持续攀升，散热成为关键瓶颈。"),
        ("方案核心", "几何耦合 + 高通量均热 + 系统级协同。3D 共形封装贴合异形芯片，TC ≥ 680 W/m·K。"),
        ("客户验证进展", "计划 Q3 完成柔性垫片工程样件，送测 AI 芯片厂商。"),
    ]
    y_off = 0
    for h, c in blocks:
        p = tf.paragraphs[0] if y_off == 0 else tf.add_paragraph()
        p.text = h
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_HIGHLIGHT
        if y_off > 0:
            p.space_before = Pt(12)
        p2 = tf.add_paragraph()
        p2.text = c
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(4)
        y_off += 1
    p3 = tf.add_paragraph()
    p3.text = "应用场景：AI 芯片、工业激光器、超算中心、IGBT 模块、5G 基站、相控阵雷达。"
    p3.font.size = Pt(9)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(8)
    # 客户收益对比表（全宽，位于页面下方）
    add_line(slide, 0.8, 4.4, 12.5, 4.4)
    benefit_rows = [
        ["指标", "传统方案", "ToSpike方案", "客户收益"],
        ["芯片结温", "95℃", "83℃", "降低12℃"],
        ["散热器体积", "100%", "70%", "减重30%"],
        ["界面热阻", "0.15 K/W", "0.09 K/W", "降低40%"],
    ]
    bw, bh = [2.2, 2.8, 2.8, 3.2], 0.42
    by, bx = 4.5, 0.8
    for ri, row in enumerate(benefit_rows):
        if ri > 0:
            add_line(slide, bx, by + ri * bh - 0.02, bx + sum(bw), by + ri * bh - 0.02)
        for ci, cell in enumerate(row):
            if ci >= len(bw):
                break
            x = bx + sum(bw[:ci])
            box = slide.shapes.add_textbox(Inches(x), Inches(by + ri * bh), Inches(bw[ci] - 0.05), Inches(bh - 0.04))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(11)
            if ri == 0:
                p.font.bold = True
                p.font.color.rgb = ACCENT_BLUE
            elif ci == 3 and ri > 0:  # 客户收益列高亮（降低12℃、减重30%、降低40%）
                p.font.bold = True
                p.font.color.rgb = ACCENT_HIGHLIGHT
                set_mono_font(p, 11)
            else:
                p.font.color.rgb = TEXT_BODY if ci > 0 else TEXT_MUTED
    return slide


def add_slide_7_roadmap(prs):
    """第五部分：三级火箭商业模型图（Mermaid 关系图风格，P0→P1→P2 层级递进）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide, "CURRENT\nPROGRESS")
    add_title_box(slide, "路线图与当前进展 — 「三级火箭」业务优先级与估值路径", y=0.55, font_size=20)
    add_line(slide, 0.8, 1.65, 12.5, 1.65)
    # 三级火箭关系图（P0→P1→P2 层级递进）
    rocket_path = gen_three_stage_rocket_img()
    if rocket_path and os.path.exists(rocket_path):
        slide.shapes.add_picture(rocket_path, Inches(1.5), Inches(1.9), Inches(10.3), Inches(2.2))
    # 底部产品堆栈摘要
    add_line(slide, 0.8, 4.3, 12.5, 4.3)
    detail = slide.shapes.add_textbox(Inches(0.8), Inches(4.45), Inches(11.7), Inches(1.8))
    tf = detail.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "P0 生存：大功率无线线圈（合同已到位）、金刚石铜（TC≥680，Q2 目标 TC>1000 认证）"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_BODY
    p2 = tf.add_paragraph()
    p2.text = "P1 估值：柔性垫片（Q3 工程样件）、全碳平台（PCT 准备中）、金刚石不粘锅（Q4 品牌联名）"
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_BODY
    p2.space_before = Pt(4)
    p3 = tf.add_paragraph()
    p3.text = "P2 护城河：固态电池骨架（模量>30 GPa 验证）、水处理电极（维持专利）"
    p3.font.size = Pt(10)
    p3.font.color.rgb = TEXT_BODY
    p3.space_before = Pt(4)
    # 底部总结
    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5))
    tf = foot.text_frame
    p = tf.paragraphs[0]
    p.text = "三级火箭模式：P0 保生存，P1 定估值，P2 筑护城河"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_HIGHLIGHT
    p.alignment = PP_ALIGN.CENTER
    return slide


def add_slide_fabless_blackbox(prs):
    """生产模式页：Fabless 黑盒策略，突出核心技术留在公司内部"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, "生产模式 — Fabless 黑盒策略：IP 安全防护", y=0.55, font_size=22)
    add_line(slide, 0.8, 1.65, 12.5, 1.65)
    # 黑盒工序流程图
    flow_path = gen_blackbox_flowchart_img()
    if flow_path and os.path.exists(flow_path):
        slide.shapes.add_picture(flow_path, Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.5))
    # 说明文字
    desc = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.5))
    tf = desc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "黑盒策略核心"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p2 = tf.add_paragraph()
    p2.text = "核心粉末改性、专利包衣处理等关键工序严格保留在簇锋内部实验室，仅输出改性母料至外协代工厂。真空钎焊、成型、机加工等通用工序外协，成品送回内部质检。核心技术不可逆向，IP 全链条可控。"
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_BODY
    p2.space_before = Pt(4)
    return slide


def add_slide_8_team(prs):
    """团队与治理：29项专利/17年军工经验加粗 + 治理机制卡片 + 团队合影占位"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, "专业互补的铁三角，清晰稳固的治理结构", y=0.55)
    add_line(slide, 0.8, 1.65, 12.5, 1.65)
    add_line(slide, 6.8, 1.8, 6.8, 6.2)
    # 左侧：核心团队（关键数据加粗）
    team_items = [
        ("王波（技术研发与应用负责人）", "工学博士，累计申请并获得授权发明专利及实用新型专利达 ", "29 项", "。擅长针对钛合金、镍基高温合金、碳化硅及 CFRP 等极端工况，通过 sp²–sp³ 界面参数的深度匹配提供定制化方案。"),
        ("何军（工程设计与结构负责人）", "机械工程硕士，高级工程师，", "17 年以上军工电子设备研发经验", "。擅长复杂电子设备的结构-热耦合一体化设计，具备从系统方案到工程实现的全链路能力。"),
        ("张工（应用工程与现场支持负责人）", "资深金刚石工具应用专家，深耕超硬磨料磨具及材料切削机理。擅长攻克复杂工业磨削难题，致力于高端新能源载具装配工具的创新应用与市场推广。", None, None),
    ]
    y = 1.9
    for name_role, pre, bold_part, post in team_items:
        nbox = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(5.8), Inches(0.35))
        tf = nbox.text_frame
        p = tf.paragraphs[0]
        p.text = name_role
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_HIGHLIGHT
        y += 0.42
        cbox = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(5.8), Inches(0.9))
        tf = cbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if bold_part:
            p.clear()
            r1 = p.add_run()
            r1.text = pre
            r1.font.size = Pt(10)
            r1.font.color.rgb = TEXT_BODY
            r2 = p.add_run()
            r2.text = bold_part
            r2.font.bold = True
            r2.font.size = Pt(10)
            r2.font.color.rgb = ACCENT_HIGHLIGHT
            r3 = p.add_run()
            r3.text = post or ""
            r3.font.size = Pt(10)
            r3.font.color.rgb = TEXT_BODY
        else:
            p.text = pre
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_BODY
        y += 1.0
    # 右侧：治理机制卡片 + 团队合影占位
    add_rect_card(slide, 7.0, 1.9, 5.2, 2.2)
    gov = slide.shapes.add_textbox(Inches(7.15), Inches(2.0), Inches(4.9), Inches(2.0))
    tf = gov.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "治理机制"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_HIGHLIGHT
    p2 = tf.add_paragraph()
    p2.text = "《创始合伙协议》：4年股权成熟期、分级回购、递延薪酬。"
    p2.font.size = Pt(9)
    p2.font.color.rgb = TEXT_BODY
    p2.space_before = Pt(4)
    p3 = tf.add_paragraph()
    p3.text = "《一致行动协议》：核心决策效率，创始团队稳固控制权。"
    p3.font.size = Pt(9)
    p3.font.color.rgb = TEXT_BODY
    p3.space_before = Pt(2)
    p4 = tf.add_paragraph()
    p4.text = "10 项核心专利布局"
    p4.font.size = Pt(10)
    p4.font.bold = True
    p4.font.color.rgb = ACCENT_HIGHLIGHT
    p4.space_before = Pt(4)
    add_rect_card(slide, 7.0, 4.3, 5.2, 1.8)
    ph = slide.shapes.add_textbox(Inches(7.15), Inches(4.7), Inches(4.9), Inches(1.2))
    tf = ph.text_frame
    p = tf.paragraphs[0]
    p.text = "【插入：团队合影】"
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER
    return slide


def add_slide_9a_funding_overview(prs):
    """融资需求概览：估值逻辑 + 核心条款"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, "融资需求——加速产品验证与平台布局", y=0.55)
    add_line(slide, 0.8, 1.65, 12.5, 1.65)
    # 三列卡片
    add_rect_card(slide, 0.8, 1.9, 3.8, 1.6)
    c1 = slide.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(3.5), Inches(1.35))
    tf = c1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "融资金额"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_HIGHLIGHT
    p2 = tf.add_paragraph()
    p2.text = "800万 - 1000万人民币"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_HIGHLIGHT
    p2.space_before = Pt(6)
    add_rect_card(slide, 4.8, 1.9, 3.8, 1.6)
    c2 = slide.shapes.add_textbox(Inches(4.95), Inches(2.0), Inches(3.5), Inches(1.35))
    tf = c2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "出让股份"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_HIGHLIGHT
    p2 = tf.add_paragraph()
    p2.text = "10% - 12%（投后估值约8000万 - 1亿）"
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_HIGHLIGHT
    p2.space_before = Pt(6)
    add_rect_card(slide, 8.8, 1.9, 3.8, 1.6)
    c3 = slide.shapes.add_textbox(Inches(8.95), Inches(2.0), Inches(3.5), Inches(1.35))
    tf = c3.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "估值逻辑"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_HIGHLIGHT
    p2 = tf.add_paragraph()
    p2.text = "对标一盛新材料等热管理材料企业"
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_BODY
    p2.space_before = Pt(6)
    # 核心逻辑
    logic = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.7), Inches(1.2))
    tf = logic.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "核心逻辑"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p2 = tf.add_paragraph()
    p2.text = "0→1 物理闭环已完成，本轮资金用于 1→10 标准化测试与 Tier-1 供应链准入准备（如 NVIDIA、华为）。"
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_BODY
    p2.space_before = Pt(4)
    return slide


def add_rect_red_box(slide, left, top, width, height):
    """红框：融资第一优先级等强调"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xfe, 0xe2, 0xe2)  # 浅红底
    shape.line.color.rgb = ACCENT_RED
    shape.line.width = Pt(1.5)
    shape.adjustments[0] = 0.03
    return shape


def add_slide_9b_funding_details(prs):
    """资金用途与下一轮里程碑：400万拆细 + 红框标出第三方认证"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_svg_top_bar(slide, height_inch=0.35)
    add_pdf_header_footer(slide)
    add_title_box(slide, "资金用途与下一轮里程碑", y=0.55)
    add_line(slide, 0.8, 1.65, 12.5, 1.65)
    # 产品研发与验证（400万）拆细
    add_line(slide, 0.8, 1.9, 12.5, 1.9)
    h1 = slide.shapes.add_textbox(Inches(0.8), Inches(1.95), Inches(11.7), Inches(0.4))
    tf = h1.text_frame
    p = tf.paragraphs[0]
    p.text = "产品研发与验证（400万）"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    # 第三方权威认证（150万）——红框
    add_rect_red_box(slide, 0.8, 2.4, 11.7, 1.1)
    r1 = slide.shapes.add_textbox(Inches(0.95), Inches(2.5), Inches(11.4), Inches(0.9))
    tf = r1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "第三方权威认证（150万）"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED
    p2 = tf.add_paragraph()
    p2.text = "融资后第一优先级投入，计划 Q2 完成 TC>1000 认证"
    p2.font.size = Pt(10)
    p2.font.italic = True
    p2.font.color.rgb = ACCENT_RED
    p2.space_before = Pt(4)
    # 中试线、工程样件
    items = [
        "中试线建设（200万）：高性能碳基复合材料中试生产线（Pilot Line）建设。",
        "工程样件送测（50万）：柔性垫片工程样件送测、固态电池骨架材料预研。",
    ]
    y = 3.6
    for item in items:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.7), Inches(0.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "• " + item
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_BODY
        y += 0.55
    # 其他用途
    add_line(slide, 0.8, 4.7, 12.5, 4.7)
    other = [
        "专利与IP布局（200万）：核心专利 PCT 国际申请、专利维护。",
        "团队建设与运营（300万）：核心团队薪酬（含递延薪酬清偿）、招聘研发人员、18 个月运营开支。",
        "市场与销售（100万）：行业展会、客户拜访、样品制作。",
    ]
    y = 4.75
    for item in other:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.7), Inches(0.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "• " + item
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_BODY
        y += 0.5
    # 下一轮里程碑
    add_line(slide, 0.8, 6.3, 12.5, 6.3)
    mh = slide.shapes.add_textbox(Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.4))
    tf = mh.text_frame
    p = tf.paragraphs[0]
    p.text = "下一轮里程碑（Pre-A 轮）——做到什么，估值到多少"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_HIGHLIGHT
    m1 = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5))
    tf = m1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• 至少 1 款产品获得头部客户（AI 芯片/军工/新能源）验证意向 → 估值预期 1.5-2 亿。"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_BODY
    p2 = tf.add_paragraph()
    p2.text = "• 完成核心技术的 PCT 国际专利布局。"
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_BODY
    p2.space_before = Pt(2)
    p3 = tf.add_paragraph()
    p3.text = "• 固态电池骨架材料完成实验室数据验证。"
    p3.font.size = Pt(11)
    p3.font.color.rgb = TEXT_BODY
    p3.space_before = Pt(2)
    return slide


# 合作愿景页专用色（深色背景 + 主色#00B4D8）
BG_VISION = RGBColor(0x02, 0x3b, 0x99)      # #023B99 与 SVG 主色一致
TEXT_VISION = RGBColor(0xf5, 0xf5, 0xf5)     # 浅灰白
ACCENT_VISION = RGBColor(0x73, 0xdb, 0xff)   # #73DBFF 与 SVG 强调色一致
LINE_VISION = RGBColor(0x33, 0x33, 0x33)     # 浅色分隔线


def add_slide_vision_merged(prs):
    """合作愿景与核心理念（合并页）：深色背景 + 科技感四角边框"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_VISION)
    add_tech_corner_frame(slide, color=ACCENT_VISION)
    add_footer(slide, is_dark=True)
    # 标题：合作愿景与核心理念（32pt，#00B4D8）
    title_box = slide.shapes.add_textbox(Inches(MARGIN), Inches(0.6), Inches(13.333 - 2 * MARGIN), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "合作愿景与核心理念"
    apply_title_font(p, 32)
    p.font.color.rgb = ACCENT_VISION
    p.alignment = PP_ALIGN.CENTER
    # 上部：我们寻找的伙伴（约占2/3，2x2网格）
    section_title = slide.shapes.add_textbox(Inches(MARGIN), Inches(1.6), Inches(11.2), Inches(0.6))
    tf = section_title.text_frame
    p = tf.paragraphs[0]
    p.text = "我们寻找的伙伴"
    apply_title_font(p, 24)
    p.font.color.rgb = ACCENT_VISION
    partners = [
        ("研究机构", "对界面物理、材料底层具有长期研究兴趣"),
        ("产业方", "愿意在应用场景中共同定义产品边界"),
        ("技术团队", "能够理解「平台价值」而非单点产品"),
        ("战略投资者", "关注中长期技术拐点、具备耐心与独立判断能力"),
    ]
    col_w = 5.4
    row_h = 1.4
    for i, (ptype, pdesc) in enumerate(partners):
        col, row = i % 2, i // 2
        x = MARGIN + col * (col_w + 0.4)
        y = 2.3 + row * row_h
        # 极简线条图标：小矩形块，主色#00B4D8
        icon_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(0.25), Inches(0.35))
        icon_shape.fill.solid()
        icon_shape.fill.fore_color.rgb = ACCENT_VISION
        icon_shape.line.width = Pt(0)
        # 内容
        content = slide.shapes.add_textbox(Inches(x + 0.35), Inches(y), Inches(col_w - 0.45), Inches(row_h - 0.1))
        tf = content.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ptype
        apply_body_font(p, 24)
        p.font.bold = True
        p.font.color.rgb = TEXT_VISION
        try:
            p.paragraph_format.line_spacing = 1.4
        except Exception:
            pass
        p2 = tf.add_paragraph()
        p2.text = pdesc
        apply_body_font(p2, 20)  # 正文20pt
        p2.font.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
        p2.space_before = Pt(4)
        try:
            p2.paragraph_format.line_spacing = 1.4
        except Exception:
            pass
    # 分隔线（约在页面2/3处，7.5*2/3≈5.0）
    sep = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(MARGIN), Inches(5.0), Inches(13.333 - MARGIN), Inches(5.0))
    sep.line.color.rgb = LINE_VISION
    sep.line.width = Pt(0.5)
    # 下部：核心理念（约占1/3）
    quote_box = slide.shapes.add_textbox(Inches(MARGIN + 0.5), Inches(5.3), Inches(12.2), Inches(1.2))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "❝ 技术的真正价值，往往并不在它被命名的那一刻，而在它被理解之前，所保留的那段空间。 ❞"
    apply_body_font(p, 28)
    p.font.italic = True
    p.font.color.rgb = TEXT_VISION
    p.alignment = PP_ALIGN.CENTER
    try:
        p.paragraph_format.line_spacing = 1.4
    except Exception:
        pass
    sub_quote = slide.shapes.add_textbox(Inches(MARGIN), Inches(6.55), Inches(13.333 - 2 * MARGIN), Inches(0.55))
    tf = sub_quote.text_frame
    p = tf.paragraphs[0]
    p.text = "我们是一个不断演化的材料与物理平台，期待认知层面一致的伙伴。"
    apply_body_font(p, 18)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER
    return slide


def main():
    # 确保 SVG 资源已提取（LOGO、产品图等，来自 20260305 公司宣传ppt 2纯净版.svg）
    if not os.path.exists(SVG_LOGO_PATH):
        try:
            import sys
            _script_dir = os.path.dirname(os.path.abspath(__file__))
            sys.path.insert(0, _script_dir)
            import extract_svg_images
            extract_svg_images.extract_images()
        except Exception:
            pass
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    add_cover_slide(prs)
    add_slide_2_opening(prs)
    add_slide_3_company(prs)
    add_slide_4_stages(prs)
    add_slide_5_principle(prs)
    add_slide_6_tech_compare(prs)
    add_slide_7_diamond_copper(prs)
    add_slide_8_patent(prs)
    add_slide_9_placeholder(prs)
    add_slide_10_rocket(prs)
    add_slide_11_team(prs)
    add_slide_12_funding(prs)
    add_slide_13_funding_details(prs)
    add_slide_14_vision(prs)

    out_dir = os.path.join(os.path.dirname(__file__), "..", "PPT files")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "簇锋科技融资路演PPT - 测试.pptx")
    prs.save(out_path)
    print(f"PPT 已生成: {os.path.abspath(out_path)}")


if __name__ == "__main__":
    main()
