#!/usr/bin/env python3
"""从 簇锋科技融资路演PPT.pptx 提取所有文字大纲"""
import os
import glob

ppt_dir = os.path.join(os.path.dirname(__file__), "..", "PPT files")
# 找原始版（不含"测试"）
pattern = os.path.join(ppt_dir, "簇锋科技融资路演PPT.pptx")
alt = os.path.join(ppt_dir, "*融资路演*.pptx")

def find_original():
    for f in os.listdir(ppt_dir):
        if f.endswith(".pptx") and "测试" not in f and "融资" in f:
            return os.path.join(ppt_dir, f)
    return None

def extract():
    from pptx import Presentation
    path = find_original()
    if not path:
        print("未找到 簇锋科技融资路演PPT.pptx")
        return []
    prs = Presentation(path)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    row_texts = []
                    for cell in row.cells:
                        t = cell.text.strip() if cell.text else ""
                        if t:
                            row_texts.append(t)
                    if row_texts:
                        texts.append(" | ".join(row_texts))
            elif hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        slides_text.append({"idx": i+1, "texts": texts})
    return slides_text

if __name__ == "__main__":
    out_path = os.path.join(ppt_dir, "ppt_text_outline.txt")
    with open(out_path, "w", encoding="utf-8") as f:
        for s in extract():
            f.write(f"=== Slide {s['idx']} ===\n")
            for t in s["texts"]:
                f.write(t + "\n")
            f.write("\n")
    print(f"已保存到 {out_path}")
