#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
对比生成的 PPT 与 SVG 文稿一致性
从 PPT 和 SVG 提取文字，逐页对比并报告差异
"""
import os
import re
import sys

# Windows 控制台 UTF-8 输出
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8")

PPT_DIR = os.path.join(os.path.dirname(__file__), "..", "PPT files")
SVG_DIR = os.path.join(PPT_DIR, "svg_output")

def find_test_pptx():
    for f in os.listdir(PPT_DIR):
        if f.endswith(".pptx") and "测试" in f:
            return os.path.join(PPT_DIR, f)
    return None

def extract_ppt_text(pptx_path):
    from pptx import Presentation
    prs = Presentation(pptx_path)
    slides = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    row_texts = [c.text.strip() for c in row.cells if c.text]
                    if row_texts:
                        texts.append(" | ".join(row_texts))
            elif hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        slides.append(texts)
    return slides

def extract_svg_text(svg_path):
    """从 SVG 提取 tspan/text 内容"""
    with open(svg_path, "r", encoding="utf-8") as f:
        content = f.read()
    # 匹配 <tspan> 或 <text> 内的文本，排除 Aspose 水印
    texts = re.findall(r'<tspan[^>]*>([^<]+)</tspan>', content)
    texts += re.findall(r'<text[^>]*>([^<]+)</text>', content)
    # 过滤水印和空行
    skip = ("Evaluation only", "Aspose", "Copyright", "Created with")
    return [t.strip() for t in texts if t.strip() and not any(s in t for s in skip)]

def _normalize(s):
    """去除空格、换行等便于对比"""
    return "".join(s.split()).replace("&amp;", "&")

def compare():
    pptx = find_test_pptx()
    if not pptx or not os.path.exists(pptx):
        print("未找到 PPT 文件")
        return False
    if not os.path.isdir(SVG_DIR):
        print("未找到 svg_output 目录，请先运行 pptx_to_svg.py")
        return False
    ppt_slides = extract_ppt_text(pptx)
    base = "簇锋科技融资路演PPT - 测试"
    all_ok = True
    for i in range(len(ppt_slides)):
        svg_path = os.path.join(SVG_DIR, f"{base}_slide_{i+1:02d}.svg")
        if not os.path.exists(svg_path):
            print(f"Slide {i+1}: 缺少 SVG 文件")
            all_ok = False
            continue
        svg_texts = extract_svg_text(svg_path)
        svg_combined = _normalize("".join(svg_texts))
        ppt_texts = ppt_slides[i]
        missing = []
        for t in ppt_texts:
            if len(t) <= 2:
                continue
            nt = _normalize(t)
            if nt not in svg_combined:
                missing.append(t)
        if missing:
            print(f"Slide {i+1}: 以下内容在 SVG 中缺失或未完整呈现:")
            for m in missing:
                preview = m[:80] + ("..." if len(m) > 80 else "")
                print(f"  - {preview}")
            all_ok = False
    if all_ok:
        print("对比完成：PPT 与 SVG 文稿一致")
    return all_ok

if __name__ == "__main__":
    compare()
